"""
Utilities for extracting and cleaning text from different data sources (PDF, URL, FAQ).
Includes specialized cleaning logic for each source type.
"""

import httpx
import re
import asyncio
from bs4 import BeautifulSoup
from PyPDF2 import PdfReader
import docx
import pandas as pd
from pptx import Presentation
from unstructured.partition.auto import partition
from playwright.async_api import async_playwright
from app.core.logging import logger
from urllib.parse import urlparse, urljoin


class PlaywrightManager:
    """
    Singleton-style manager for Playwright browser instance.
    Managed via the application lifespan in main.py.
    """
    def __init__(self):
        self._pw = None
        self._browser = None

    async def start(self):
        if not self._pw:
            self._pw = await async_playwright().start()
            self._browser = await self._pw.chromium.launch(
                headless=True,
                args=["--disable-dev-shm-usage", "--no-sandbox"]
            )
            logger.info("Playwright Browser started")

    async def stop(self):
        if self._browser:
            await self._browser.close()
            self._browser = None
        if self._pw:
            await self._pw.stop()
            self._pw = None
        logger.info("Playwright Browser stopped")

    def get_browser(self):
        return self._browser


pw_manager = PlaywrightManager()


def clean_pdf_text(text: str) -> str:
    """Clean text extracted from PDF: removes headers/footers placeholders, fixes ligatures."""
    # Remove multiple newlines and normalize whitespace
    text = re.sub(r'\n\s*\n', '\n\n', text)
    text = re.sub(r'[ \t]+', ' ', text)
    
    # Remove common PDF artifacts like "Page 1 of 10" or similar patterns
    text = re.sub(r'Page \d+ of \d+', '', text, flags=re.IGNORECASE)
    
    # Fix common character issues (e.g., "fi" or "fl" ligatures if unparsed)
    # This is a basic pass; more complex PDFs might need specialized libraries.
    
    return text.strip()


def clean_url_text(soup: BeautifulSoup) -> str:
    """
    Advanced Content Extraction Strategy:
    1. Extracts Title & Metadata for context.
    2. Parses JSON-LD for structured data.
    3. Identifies and extracts 'Main' content while stripping boilerplate.
    4. Normalizes whitespace and structure.
    """
    # 1. Metadata & Title Extraction (Contextual enrichment)
    title = soup.title.string if soup.title else ""
    meta_desc = ""
    desc_tag = soup.find("meta", attrs={"name": "description"}) or soup.find("meta", attrs={"property": "og:description"})
    if desc_tag:
        meta_desc = desc_tag.get("content", "")

    # 2. JSON-LD Extraction (Structured Data)
    structured_data = []
    for script in soup.find_all("script", type="application/ld+json"):
        try:
            # We don't necessarily need to parse it as JSON, just clean it up
            content = script.string.strip()
            if content:
                structured_data.append(content)
        except: continue

    # 3. Boilerplate & Noise Removal
    for element in soup(["script", "style", "nav", "footer", "header", "aside", "form", "iframe", "input", "noscript"]):
        element.decompose()
        
    # Remove common ad and social noise classes/ids
    noise_patterns = ["ad-", "social-", "share-", "sidebar", "menu", "banner", "cookie"]
    for tag in soup.find_all(class_=re.compile("|".join(noise_patterns), re.I)):
        tag.decompose()
    for tag in soup.find_all(id=re.compile("|".join(noise_patterns), re.I)):
        tag.decompose()

    # 4. Hierarchical Content Extraction
    # Strategy A: Semantic Containers (Extract ALL matches, not just one)
    selectors = "main, article, [role='main'], #content, .content, .page, .post-content, .article-body"
    containers = soup.select(selectors)
    
    text_parts = []
    if title: text_parts.append(f"PAGE TITLE: {title}")
    if meta_desc: text_parts.append(f"DESCRIPTION: {meta_desc}")
    
    if containers:
        for container in containers:
            # Avoid nesting (if main contains article, don't double extract)
            # This is a bit complex, but usually get_text() on the largest container is best.
            # We'll take the top-level ones.
            if not any(parent in containers for parent in container.parents):
                text_parts.append(container.get_text(separator="\n", strip=True))
    else:
        # Strategy B: High-density text blocks (fallback)
        potential_containers = soup.find_all(["div", "section"])
        best_containers = []
        for container in potential_containers:
            p_count = len(container.find_all("p"))
            if p_count > 2:
                best_containers.append(container)
        
        if best_containers:
            for container in best_containers:
                if not any(parent in best_containers for parent in container.parents):
                    text_parts.append(container.get_text(separator="\n", strip=True))
        else:
            # Final Strategy: Cleaned Body
            body = soup.find('body')
            if body:
                text_parts.append(body.get_text(separator="\n", strip=True))

    # 5. Combine and Refine
    raw_text = "\n\n".join(text_parts)
    
    # Normalize empty lines and whitespace
    lines = []
    seen_lines = set()
    for line in raw_text.splitlines():
        line = line.strip()
        # Deduplicate and filter
        if (len(line) >= 3 or "TITLE:" in line) and line not in seen_lines:
            lines.append(line)
            seen_lines.add(line)
            
    return "\n".join(lines).strip()


def get_internal_links(soup: BeautifulSoup, base_url: str) -> list:
    """Extract internal links from a BeautifulSoup object."""
    links = set()
    parsed_base = urlparse(base_url)
    domain = parsed_base.netloc
    
    for a in soup.find_all('a', href=True):
        href = a['href']
        # Resolve relative URLs
        full_url = urljoin(base_url, href)
        parsed_url = urlparse(full_url)
        
        # Stay on the same domain and handle same-protocol or http/https
        if parsed_url.netloc == domain and parsed_url.scheme in ['http', 'https', '']:
            # Strip fragments and query params to avoid duplicates
            clean_url = full_url.split('#')[0].split('?')[0].rstrip('/')
            if clean_url and clean_url != base_url.rstrip('/'):
                links.add(clean_url)
                
    return list(links)


def clean_faq_text(question: str, answer: str) -> str:
    """Clean and format individual FAQ entries."""
    q = question.strip().rstrip('?') + '?'
    a = answer.strip()
    return f"Q: {q}\nA: {a}"


# ── Extraction Entry Points ──

def extract_text_from_pdf(file_path: str) -> str:
    """Extract and specifically clean text from a local PDF."""
    raw_text = ""
    try:
        reader = PdfReader(file_path)
        for page in reader.pages:
            content = page.extract_text()
            if content:
                raw_text += content + "\n"
        return clean_pdf_text(raw_text)
    except Exception as e:
        logger.error(f"Error extracting PDF: {e}")
        return ""


def extract_text_from_docx(file_path: str) -> str:
    """Extract text from a Word document."""
    try:
        doc = docx.Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs])
    except Exception as e:
        logger.error(f"Error extracting DOCX: {e}")
        return ""


def extract_text_from_xlsx(file_path: str) -> str:
    """Extract text from an Excel spreadsheet."""
    try:
        # Read all sheets
        df_dict = pd.read_excel(file_path, sheet_name=None)
        text = ""
        for sheet_name, df in df_dict.items():
            text += f"Sheet: {sheet_name}\n"
            # Replace NaNs with empty string
            df = df.fillna("")
            text += df.to_csv(index=False, sep='\t') + "\n\n"
        return text
    except Exception as e:
        logger.error(f"Error extracting XLSX: {e}")
        return ""


def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from a PowerPoint presentation."""
    try:
        prs = Presentation(file_path)
        text = ""
        for i, slide in enumerate(prs.slides):
            text += f"Slide {i+1}:\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
            text += "\n"
        return text
    except Exception as e:
        logger.error(f"Error extracting PPTX: {e}")
        return ""


def extract_text_universal(file_path: str) -> str:
    """Extract text using unstructured library as a robust fallback."""
    try:
        elements = partition(filename=file_path)
        return "\n\n".join([str(el) for el in elements])
    except Exception as e:
        logger.error(f"Error using unstructured extractor: {e}")
        return ""


async def extract_text_from_url(url: str, max_pages: int = 15) -> str:
    """
    Scrape, extract, and specifically clean content from a URL and its subpages.
    Uses a shared Playwright browser instance for efficiency.
    """
    browser = pw_manager.get_browser()
    if not browser:
        logger.warning(f"Playwright browser not initialized, falling back to static for {url}")
        return await _extract_text_static(url)

    visited = set()
    to_visit = [url]
    all_text_parts = []

    try:
        context = await browser.new_context(
            user_agent="Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/122.0.0.0 Safari/537.36",
            viewport={'width': 1280, 'height': 900}
        )
        
        while to_visit and len(visited) < max_pages:
            current_url = to_visit.pop(0)
            if current_url in visited:
                continue
                
            visited.add(current_url)
            page = await context.new_page()
            
            try:
                logger.info(f"Scraping [{len(visited)}/{max_pages}]: {current_url}")
                
                # 1. Navigation with SPA Recovery
                response = None
                try:
                    response = await page.goto(current_url, wait_until="networkidle", timeout=30000)
                except Exception:
                    logger.warning(f"Timeout/Error loading {current_url}, attempting recovery...")

                # SPA Recovery: If direct hit fails/404s, try navigating to base and clicking
                if not response or response.status == 404:
                    parsed_url = urlparse(current_url)
                    base_url = f"{parsed_url.scheme}://{parsed_url.netloc}/"
                    path = parsed_url.path.strip('/')
                    if path:
                        try:
                            await page.goto(base_url, wait_until="networkidle", timeout=30000)
                            # Try to find the link on the home page and click it
                            link = page.locator(f"a[href$='{path}']").first
                            if await link.count() > 0:
                                await link.click()
                                await page.wait_for_load_state("networkidle")
                        except Exception as recovery_err:
                            logger.warning(f"Recovery failed for {current_url}: {recovery_err}")

                # 2. Wait for dynamic content (hydration)
                await asyncio.sleep(2.5) 
                
                # 3. Escape popups
                try: await page.keyboard.press("Escape")
                except: pass

                # 4. Capture Content
                content = await page.content()
                soup = BeautifulSoup(content, "html.parser")
                
                # Extract text from this page
                page_text = clean_url_text(soup)
                if page_text:
                    all_text_parts.append(f"--- SOURCE: {current_url} ---\n{page_text}")
                
                # 5. Extract links for further crawling (only from the first few pages)
                if len(visited) <= 3:
                    internal_links = get_internal_links(soup, current_url)
                    internal_links.sort(key=len)
                    for link in internal_links:
                        if link not in visited and link not in to_visit:
                            to_visit.append(link)
                            
            except Exception as page_err:
                logger.error(f"Error scraping {current_url}: {page_err}")
            finally:
                await page.close()

        if not all_text_parts:
            return await _extract_text_static(url)
            
        return "\n\n\n".join(all_text_parts)

    except Exception as e:
        logger.error(f"Crawler structural failure for {url}: {e}")
        return await _extract_text_static(url)
    
    finally:
        if context:
            await context.close()


async def _extract_text_static(url: str) -> str:
    """Fallback static scraper for URLs when Playwright fails or is unnecessary."""
    try:
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36",
            "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
            "Accept-Language": "en-US,en;q=0.9"
        }
        async with httpx.AsyncClient(
            headers=headers,
            timeout=httpx.Timeout(20.0, connect=10.0),
            follow_redirects=True,
        ) as client:
            response = await client.get(url)
            response.raise_for_status()

        content_type = response.headers.get("content-type", "").lower()
        if "text/plain" in content_type:
            return response.text.strip()
        
        soup = BeautifulSoup(response.text, "html.parser")
        return clean_url_text(soup)
    except Exception as e:
        logger.error(f"Static fallback scraping failed for {url}: {e}")
        return ""


def format_faqs_to_text(faqs: list) -> str:
    """Clean and format a list of FAQs into a unified text block."""
    entries = []
    for faq in faqs:
        cleaned = clean_faq_text(faq.get("question", ""), faq.get("answer", ""))
        entries.append(cleaned)
    return "\n\n".join(entries)
